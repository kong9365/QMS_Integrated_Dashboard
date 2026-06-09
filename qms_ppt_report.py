# -*- coding: utf-8 -*-
"""qms_ppt_report.py — OOS 마감회의 & GMP 보고서 PowerPoint(.pptx) 생성기 (KD-MoaQ).

설계 원칙
- **신규 순수함수만 추가**. 기존 표시 함수 ``qms_oos_dashboard_panels.render_oos_gmp``
  및 도메인/건수기여도 계산 로직은 일절 수정하지 않는다(탭 레이아웃·디자인 불변).
- render_oos_gmp 의 집계(건수기여도 기반 groupby/sum)를 **그대로 미러링**하여, 화면의
  표/차트를 **python-pptx 네이티브 요소(표·차트)**로 충실히 재현한다.
  → 슬라이드 구성은 화면의 2행 레이아웃을 따른다(1행=현황표·월별·Analyst error,
     2행=원인 대분류·소분류·시험종류별).
- 차트는 python-pptx 네이티브 차트(편집 가능) → kaleido(이미지 변환) 불필요.

공개 API
- build_oos_gmp_report_pptx(filtered, df_full, primary_year, prev_year, year_col,
      month_col, safe_pct, completed_keywords, *,
      as_of=None, project_label="OOS (Out of Specification)", filter_note="") -> bytes
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Optional, Tuple

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

_MONTH_LABELS = [f"{m}월" for m in range(1, 13)]
_KD_RED = RGBColor(0xE8, 0x30, 0x08)     # 강조/경고
_NAVY = RGBColor(0x1F, 0x3A, 0x5F)       # 헤더/주계열
_BAR = RGBColor(0x4C, 0x78, 0xA8)        # (구) 막대색 — 하위호환
_DARK = RGBColor(0x59, 0x59, 0x59)
_GRAY = RGBColor(0x80, 0x80, 0x80)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GREEN = RGBColor(0x1F, 0x9D, 0x63)      # 완료(성공)
_INK = RGBColor(0x1B, 0x23, 0x30)        # 본문 잉크
_LINE = RGBColor(0xE1, 0xE6, 0xEF)       # 옅은선/격자
_ZEBRA = RGBColor(0xF4, 0xF7, 0xFA)
_KD_PALETTE = [_NAVY, _KD_RED, RGBColor(0x2F, 0x6F, 0xED),
               RGBColor(0x0E, 0x9A, 0xA7), RGBColor(0xE8, 0x83, 0x0C), RGBColor(0x7A, 0x5A, 0xF0)]


# ── render_oos_gmp 와 동일한 보조 계산(미러링) ──
def _monthly_counts(data: pd.DataFrame, year, year_col: str, month_col: str = "월"):
    if data is None or data.empty or year_col not in data.columns or month_col not in data.columns:
        return [0] * 12
    yr = data[data[year_col] == year]
    if yr.empty:
        return [0] * 12
    if "건수기여도" in yr.columns:
        monthly = yr.groupby(month_col, dropna=False)["건수기여도"].sum()
    else:
        monthly = yr.groupby(month_col, dropna=False).size().astype(float)
    monthly = monthly.reindex(range(1, 13), fill_value=0)
    return [round(float(x)) for x in monthly.tolist()]


def _completed_mask(df: pd.DataFrame, kw: Tuple[str, ...]) -> pd.Series:
    if "진행상태" in df.columns:
        return df["진행상태"].str.contains("|".join(kw), case=False, na=False)
    if "완료여부" in df.columns:
        return df["완료여부"] == "C"
    return pd.Series(False, index=df.index)


# ── PPT 요소 헬퍼 ──
def _textbox(slide, text, *, size=20, color=_NAVY, bold=True, left=0.5, top=0.3,
             width=12.3, height=0.7, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "맑은 고딕"
    return tb


def _panel_title(slide, text, *, top=0.35):
    """슬라이드 패널 제목(네이비 굵게) + 하단 KD Red 얇은 바."""
    _textbox(slide, text, size=22, color=_NAVY, bold=True, left=0.6, top=top, width=12.1, height=0.6)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(top + 0.62), Inches(2.4), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _KD_RED
    bar.line.fill.background()
    bar.shadow.inherit = False


def _caption(slide, text, *, left, top, width=5.8):
    """차트 위 작은 패널 캡션(네이비 굵게) — 차트 내부 제목 대체(중복 방지)."""
    _textbox(slide, text, size=13, color=_NAVY, bold=True, left=left, top=top, width=width, height=0.35)


def _footer(slide, *, page=None, as_of=""):
    """하단 푸터(브랜드 + 기준일) + 우측 페이지 번호."""
    _textbox(slide, f"KD-MoaQ · 광동제약 품질부문{('  ·  ' + as_of) if as_of else ''}",
             size=9, color=_GRAY, bold=False, left=0.6, top=7.05, width=10.5, height=0.35)
    if page is not None:
        _textbox(slide, str(page), size=9, color=_GRAY, bold=False,
                 left=12.4, top=7.05, width=0.5, height=0.35, align=PP_ALIGN.RIGHT)


def _kpi_cards(slide, cards, *, left, top, width, height):
    """요약 KPI 카드(둥근 사각형: 큰 숫자+라벨, 상단 색 보더). cards=[(label, value, accent_rgb), ...]."""
    n = len(cards)
    gap = 0.22
    cw = (width - gap * (n - 1)) / n
    for i, (label, value, accent) in enumerate(cards):
        x = left + i * (cw + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(cw), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = _WHITE
        card.line.color.rgb = _LINE
        card.line.width = Pt(1.0)
        card.shadow.inherit = False
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(cw), Inches(0.09))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = str(value)
        r1.font.size = Pt(26)
        r1.font.bold = True
        r1.font.color.rgb = _NAVY
        r1.font.name = "맑은 고딕"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = str(label)
        r2.font.size = Pt(12)
        r2.font.color.rgb = _DARK
        r2.font.name = "맑은 고딕"


def _bar_chart(slide, categories, values, *, left, top, width, height, horizontal=False):
    """막대 차트(네이티브) — 주계열 네이비·최댓값만 KD Red·데이터라벨 12pt·옅은 격자·한글 폰트.
    제목은 차트 내부 대신 패널 캡션(_caption)으로 분리(중복 방지)."""
    cd = CategoryChartData()
    cd.categories = [str(c) for c in categories]
    vals = [float(v) for v in values]
    cd.add_series("건수", vals)
    ctype = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    frame = slide.shapes.add_chart(ctype, Inches(left), Inches(top), Inches(width), Inches(height), cd)
    chart = frame.chart
    chart.has_legend = False
    chart.has_title = False
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = "0"
        dl.number_format_is_linked = False
        dl.font.size = Pt(12)
        dl.font.name = "맑은 고딕"
        dl.font.color.rgb = _INK
        series = plot.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _NAVY
        # 최댓값 막대만 KD Red 강조(첫 최댓값 1개)
        if vals:
            mx = max(vals)
            hit = False
            for i, v in enumerate(vals):
                if not hit and v == mx and v > 0:
                    pt = series.points[i]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = _KD_RED
                    hit = True
        # 축 한글 폰트 + 옅은 격자
        for ax in (chart.category_axis, chart.value_axis):
            ax.tick_labels.font.size = Pt(10)
            ax.tick_labels.font.name = "맑은 고딕"
        chart.value_axis.has_major_gridlines = True
        gl = chart.value_axis.major_gridlines.format.line
        gl.color.rgb = _LINE
        gl.width = Pt(0.5)
    except Exception:
        pass
    return chart


def build_oos_gmp_report_pptx(
    filtered: pd.DataFrame,
    df_full: pd.DataFrame,
    primary_year: int,
    prev_year: int,
    year_col: str,
    month_col: str,
    safe_pct,
    completed_keywords: Tuple[str, ...],
    *,
    as_of: Optional[str] = None,
    project_label: str = "OOS (Out of Specification)",
    filter_note: str = "",
) -> bytes:
    """마감회의 & GMP 보고서를 .pptx 바이트로 생성(render_oos_gmp 집계·레이아웃 미러링)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    as_of = as_of or datetime.now().strftime("%Y-%m-%d %H:%M")

    # ── Slide 1: 표지 ──
    s = prs.slides.add_slide(blank)
    _textbox(s, "광동제약 품질부문 · KD-MoaQ", size=15, color=_KD_RED, top=2.2, align=PP_ALIGN.CENTER)
    _textbox(s, "OOS 마감회의 & GMP 보고서", size=36, color=_NAVY, top=2.75, align=PP_ALIGN.CENTER)
    _textbox(s, f"대상: {project_label}    기준일: {as_of}", size=15, color=_GRAY, bold=False, top=3.85, align=PP_ALIGN.CENTER)
    if filter_note:
        _textbox(s, f"필터: {filter_note}", size=12, color=_GRAY, bold=False, top=4.3, align=PP_ALIGN.CENTER)

    yc = year_col if (filtered is not None and year_col in getattr(filtered, "columns", [])) else "연도"
    mc = month_col if (filtered is not None and month_col in getattr(filtered, "columns", [])) else "월"

    if filtered is None or getattr(filtered, "empty", True) or "건수기여도" not in filtered.columns:
        s2 = prs.slides.add_slide(blank)
        _textbox(s2, "표시할 OOS 데이터가 없습니다 (필터 결과 0건 또는 건수기여도 컬럼 없음).", size=18, top=3.0)
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── 집계(render_oos_gmp 미러링) ──
    total_events = round(float(filtered["건수기여도"].sum()))
    completed_ev = round(float(filtered.loc[_completed_mask(filtered, completed_keywords), "건수기여도"].sum()))
    ev_comp_rate = safe_pct(completed_ev, total_events)
    capa_df2 = pd.DataFrame()
    if "CAPA/Action item 필요여부" in filtered.columns:
        m2 = filtered["CAPA/Action item 필요여부"].fillna("").str.len() > 0
        no2 = filtered["CAPA/Action item 필요여부"].fillna("").str.contains("No Action", case=False, na=False)
        capa_df2 = filtered[m2 & ~no2]
    capa_cnt2 = round(float(capa_df2["건수기여도"].sum())) if (not capa_df2.empty and "건수기여도" in capa_df2.columns) else 0
    if not capa_df2.empty and "진행상태" in capa_df2.columns and "건수기여도" in capa_df2.columns:
        capa_prog2 = safe_pct(round(float(capa_df2.loc[_completed_mask(capa_df2, completed_keywords), "건수기여도"].sum())), capa_cnt2)
    else:
        capa_prog2 = 0.0

    # ── Slide 2: 개요 — KPI 카드 + 월별 추세(1메시지: 전반 현황) ──
    s = prs.slides.add_slide(blank)
    _panel_title(s, "시험실 이벤트 발생 현황")
    _kpi_cards(s, [
        ("전체 건수", f"{total_events:,}", _NAVY),
        ("완료 건수", f"{completed_ev:,}", _GREEN),
        ("완료율", f"{ev_comp_rate:.0f}%", _NAVY),
        ("미완료 건수", f"{total_events - completed_ev:,}", _KD_RED),
    ], left=0.6, top=1.4, width=12.1, height=1.5)
    _textbox(s, f"CAPA {capa_cnt2}건 · 진행률 {capa_prog2:.1f}%", size=12, color=_GRAY, bold=False,
             left=0.6, top=3.05, width=12.1, height=0.35)
    monthly = _monthly_counts(filtered, primary_year, yc, mc)
    _caption(s, "월별 시험실이벤트 건수", left=0.6, top=3.55, width=12.1)
    _bar_chart(s, _MONTH_LABELS, monthly, left=0.6, top=3.95, width=12.1, height=2.95)
    _footer(s, page=2, as_of=as_of)

    # ── Slide 3: 원인 분석 — 대분류 + 소분류(가로) (1메시지: 왜 발생했나) ──
    s = prs.slides.add_slide(blank)
    _panel_title(s, "시험실이벤트 원인 분석")
    if "확인된 이벤트 분류" in filtered.columns:
        major = (
            filtered[filtered["확인된 이벤트 분류"].notna() & (filtered["확인된 이벤트 분류"] != "")]
            .groupby("확인된 이벤트 분류")["건수기여도"].sum().round().sort_values(ascending=False).reset_index()
        )
        major.columns = ["분류", "건수"]
        if not major.empty:
            _caption(s, "원인 대분류", left=0.6, top=1.5, width=6.0)
            _bar_chart(s, major["분류"].tolist(), major["건수"].tolist(),
                       left=0.6, top=1.95, width=6.0, height=4.9)
    if "이상발생 원인" in filtered.columns:
        minor = (
            filtered[filtered["이상발생 원인"].notna() & (filtered["이상발생 원인"] != "")]
            .groupby("이상발생 원인")["건수기여도"].sum().round().sort_values(ascending=True).reset_index()
        )
        minor.columns = ["원인", "건수"]
        if not minor.empty:
            _caption(s, "원인 소분류", left=6.9, top=1.5, width=5.8)
            _bar_chart(s, minor["원인"].tolist(), minor["건수"].tolist(),
                       left=6.9, top=1.95, width=5.8, height=4.9, horizontal=True)
    _footer(s, page=3, as_of=as_of)

    # ── Slide 4: 유형·개선 — 시험종류별 + Analyst error(1메시지: 유형+개선효과) ──
    s = prs.slides.add_slide(blank)
    _panel_title(s, "시험종류 분포 · 개선 효과")
    if "시험종류" in filtered.columns:
        ttype = (
            filtered[filtered["시험종류"].notna() & (filtered["시험종류"] != "")]
            .groupby("시험종류")["건수기여도"].sum().round().sort_values(ascending=False).reset_index()
        )
        ttype.columns = ["시험종류", "건수"]
        if not ttype.empty:
            _caption(s, "시험종류별 발생 건수", left=0.6, top=1.5, width=6.0)
            _bar_chart(s, ttype["시험종류"].tolist(), ttype["건수"].tolist(),
                       left=0.6, top=1.95, width=6.0, height=4.9)
    if "이상발생 원인" in df_full.columns and "건수기여도" in df_full.columns and year_col in df_full.columns:
        ae_prev = round(float(df_full[(df_full[year_col] == prev_year) & (df_full["이상발생 원인"] == "Analyst error")]["건수기여도"].sum()))
        ae_curr = round(float(filtered[filtered["이상발생 원인"] == "Analyst error"]["건수기여도"].sum())) if "이상발생 원인" in filtered.columns else 0
        if ae_prev > 0 or ae_curr > 0:
            _caption(s, "Analyst error 감소 추이", left=6.9, top=1.5, width=5.8)
            _bar_chart(s, [str(prev_year), str(primary_year)], [ae_prev, ae_curr],
                       left=6.9, top=1.95, width=5.8, height=4.2)
            red = safe_pct(ae_prev - ae_curr, ae_prev) if ae_prev > 0 else 0.0
            _txt = f"감소율 -{red:.0f}%" if red > 0 else f"증가율 +{abs(red):.0f}%"
            _textbox(s, _txt, size=15, color=(_GREEN if red > 0 else _KD_RED), bold=True,
                     left=6.9, top=6.25, width=5.8, height=0.5, align=PP_ALIGN.CENTER)
    _footer(s, page=4, as_of=as_of)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
